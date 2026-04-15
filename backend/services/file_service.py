"""
Serviço de geração de arquivos — PDF, DOCX, XLSX, PPTX.
Portado de netlify/functions/files.ts
"""

import asyncio
import io
import json
import logging
import os
import re
from pathlib import Path
from typing import Optional, Tuple, Union

from backend.services.design_contract import CANVAS_PRESETS

logger = logging.getLogger(__name__)

_TEMPLATES_DIR = Path(__file__).parent / "pdf_templates"

# ── Tabelas de viewport para export responsivo ────────────────────────────────

_PAGE_VIEWPORTS = {
    "widescreen":       (1920, 1080),
    "a4-landscape":     (1122, 794),
    "a4-portrait":      (794, 1122),
    "letter-landscape": (1056, 816),
    "letter-portrait":  (816, 1056),
}

_RES_VIEWPORTS = {
    "hd-720":  (1280, 720),
    "hd-1080": (1920, 1080),
}

# ── Paged.js helpers ──────────────────────────────────────────────────────────

_PAGED_JS_CDN = "https://unpkg.com/pagedjs/dist/paged.polyfill.js"


def _should_use_paged_js(html: str, page_size: Optional[str]) -> bool:
    """
    Retorna True se o HTML deve ser renderizado com Paged.js.
    Aplica-se a documentos A4/Letter paginados que possuem margens, cabeçalhos
    e rodapés (gerados pelo a4_document_creator e similares).
    NÃO se aplica a apresentações (1920×1080) ou social media.
    """
    if page_size in ("a4-portrait", "a4-landscape", "letter-portrait", "letter-landscape"):
        return True
    # Detecta dimensões A4 portrait no CSS: width 794px + height ~1123px
    if re.search(r"width:\s*794px", html) and re.search(r"height:\s*1\d{3}px", html):
        return True
    return False


def _inject_paged_js(html: str, vw: int, vh: int) -> str:
    """
    Injeta Paged.js polyfill + regras CSS @page no HTML.
    Adiciona page-break-after em section.slide para garantir uma slide por página PDF.
    """
    inject = (
        "<style>\n"
        f"@page {{ size: {vw}px {vh}px; margin: 0; }}\n"
        "section.slide, .slide { page-break-after: always !important; break-after: page !important; }\n"
        "</style>\n"
        f'<script src="{_PAGED_JS_CDN}"></script>\n'
    )
    if "</head>" in html:
        return html.replace("</head>", inject + "</head>", 1)
    if "<body" in html:
        i = html.find("<body")
        j = html.find(">", i)
        return html[: j + 1] + inject + html[j + 1 :]
    return inject + html


# JS que revela um slide no fluxo natural do documento (sem position:fixed).
# Usado como fallback quando bounding_box() retorna None (slide estava oculto).
_REVEAL_SLIDE_JS = """(el) => {
    const all = document.querySelectorAll('.slide, .slide-container');
    all.forEach(s => {
        if (s !== el) {
            s.style.display = 'none';
        }
    });
    el.style.display = '';
    el.style.visibility = 'visible';
    el.style.opacity = '1';
    el.style.transform = 'none';
    el.style.position = '';
    el.classList.add('active');
    document.body.offsetHeight;
}"""


def _capture_slide(page, slide_els: list, idx: int, fmt: str = "png") -> Optional[bytes]:
    """
    Captura um slide usando element.screenshot() — faz scroll automático e
    captura as dimensões exatas do elemento, sem position:fixed e sem clip manual.
    Se o elemento estiver oculto (carrossel), revela-o no fluxo natural primeiro.
    """
    if idx >= len(slide_els):
        return None
    el = slide_els[idx]
    bbox = el.bounding_box()

    if not (bbox and bbox["width"] > 0 and bbox["height"] > 0):
        # Elemento oculto (carrossel) — revelar sem position:fixed
        page.evaluate(_REVEAL_SLIDE_JS, el)
        page.wait_for_timeout(300)
        bbox = el.bounding_box()

    if bbox and bbox["width"] > 0 and bbox["height"] > 0:
        # el.screenshot() rola o elemento para a view e captura suas dimensões exatas
        return el.screenshot(type=fmt)
    return None


async def generate_pdf_playwright(
    html_content: str,
    slide_index: Optional[int] = None,
    slide_indices: Optional[list] = None,
    page_size: Optional[str] = None,
    canvas_preset: Optional[str] = None,
    viewport_width: Optional[int] = None,
    viewport_height: Optional[int] = None,
) -> bytes:
    """
    Gera PDF renderizando o HTML via Playwright.
    Captura cada slide individualmente via el.screenshot() (igual ao PPTX) e
    compõe um PDF multi-página com reportlab. Para documentos simples (sem .slide),
    usa full-page screenshot.
    """

    def _sync_render() -> bytes:
        from playwright.sync_api import sync_playwright
        import io as _io
        from reportlab.lib.pagesizes import A4, letter, landscape as rl_landscape
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas as pdf_canvas

        inject_html = _inject_tailwind_if_needed(html_content)
        if viewport_width and viewport_height:
            vw, vh = viewport_width, viewport_height
        elif canvas_preset and canvas_preset in CANVAS_PRESETS:
            spec = CANVAS_PRESETS[canvas_preset]
            vw, vh = spec["width"], spec["height"]
        else:
            vw, vh = _PAGE_VIEWPORTS.get(page_size or "widescreen", (1920, 1080))

        # ── Per-slide screenshot path (igual ao PPTX) ──
        # Captura cada slide individualmente via _capture_slide() — imune a
        # CSS do body (flex, padding, overflow) e garante alinhamento perfeito.

        # Export CSS override: remove body flex/padding/overflow do design_contract
        inject_html = _inject_export_css(inject_html, vw)

        _PX_TO_PT = 72 / 96
        _RL_SIZES = {
            "widescreen":       (vw * _PX_TO_PT, vh * _PX_TO_PT),
            "a4-landscape":     rl_landscape(A4),
            "a4-portrait":      A4,
            "letter-landscape": rl_landscape(letter),
            "letter-portrait":  letter,
        }
        rl_size = _RL_SIZES.get(page_size or "widescreen", (vw * _PX_TO_PT, vh * _PX_TO_PT))
        page_w_pt, page_h_pt = rl_size

        screenshots: list[bytes] = []

        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page(viewport={"width": vw, "height": vh})
            try:
                page.set_content(inject_html, wait_until="networkidle", timeout=30_000)
                page.wait_for_timeout(200)

                slide_count = page.evaluate(
                    "() => document.querySelectorAll('.slide, .slide-container').length"
                )

                if slide_count and int(slide_count) > 1:
                    slide_els = page.query_selector_all('.slide, .slide-container')
                    total = int(slide_count)

                    # Determina quais slides capturar
                    if slide_indices is not None:
                        indices = sorted(set(i for i in slide_indices if 0 <= i < total))
                    elif slide_index is not None:
                        indices = [slide_index] if 0 <= slide_index < total else [0]
                    else:
                        indices = list(range(total))

                    for i in indices:
                        shot = _capture_slide(page, slide_els, i, "png")
                        if shot:
                            screenshots.append(shot)
                else:
                    # Documento simples — screenshot full-page
                    screenshots.append(page.screenshot(full_page=True, type="png"))
            finally:
                browser.close()

        # Monta PDF com uma página por screenshot
        pdf_buf = _io.BytesIO()
        c = pdf_canvas.Canvas(pdf_buf, pagesize=rl_size)
        for shot_bytes in screenshots:
            img_reader = ImageReader(_io.BytesIO(shot_bytes))
            iw, ih = img_reader.getSize()
            scale = min(page_w_pt / iw, page_h_pt / ih)
            draw_w, draw_h = iw * scale, ih * scale
            x = (page_w_pt - draw_w) / 2
            y_pos = page_h_pt - draw_h  # ancora no topo
            c.drawImage(img_reader, x, y_pos, draw_w, draw_h)
            c.showPage()
        c.save()
        return pdf_buf.getvalue()

    return await asyncio.to_thread(_sync_render)


async def generate_pdf_from_template(template_name: str, data: dict) -> bytes:
    """
    Gera PDF a partir de um template Jinja2 HTML pré-aprovado.
    O LLM fornece apenas os dados (JSON); o layout vem do template.
    """
    try:
        from jinja2 import Environment, FileSystemLoader, select_autoescape
    except ImportError:
        raise RuntimeError("Jinja2 não instalado. Execute: pip install jinja2")

    template_file = f"{template_name}.html"
    if not (_TEMPLATES_DIR / template_file).exists():
        available = [p.stem for p in _TEMPLATES_DIR.glob("*.html")]
        raise ValueError(
            f"Template '{template_name}' não encontrado. Disponíveis: {available}"
        )

    env = Environment(
        loader=FileSystemLoader(str(_TEMPLATES_DIR)),
        autoescape=select_autoescape(["html"]),
    )
    template = env.get_template(template_file)
    html_content = template.render(**data)
    return await generate_pdf_playwright(html_content)


def generate_pdf(title: str, content: str) -> bytes:
    """Gera PDF com reportlab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import cm

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                            topMargin=2*cm, bottomMargin=2*cm,
                            leftMargin=2*cm, rightMargin=2*cm)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "CustomTitle", parent=styles["Title"],
        fontSize=18, spaceAfter=20
    )
    body_style = ParagraphStyle(
        "CustomBody", parent=styles["Normal"],
        fontSize=12, leading=16
    )

    elements = [Paragraph(title, title_style), Spacer(1, 12)]

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 8))
        elif line.startswith("# "):
            elements.append(Paragraph(f"<b>{line[2:]}</b>", styles["Heading1"]))
        elif line.startswith("## "):
            elements.append(Paragraph(f"<b>{line[3:]}</b>", styles["Heading2"]))
        else:
            # Escape XML chars for reportlab
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            elements.append(Paragraph(safe, body_style))

    doc.build(elements)
    return buffer.getvalue()


def generate_docx(title: str, content: str) -> bytes:
    """Gera DOCX com python-docx."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)
    doc.add_paragraph("")  # spacer

    for line in content.split("\n"):
        trimmed = line.strip()
        if trimmed.startswith("# "):
            doc.add_heading(trimmed[2:], level=1)
        elif trimmed.startswith("## "):
            doc.add_heading(trimmed[3:], level=2)
        elif trimmed.startswith("### "):
            doc.add_heading(trimmed[4:], level=3)
        elif trimmed:
            p = doc.add_paragraph(trimmed)
            for run in p.runs:
                run.font.size = Pt(12)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def generate_xlsx(title: str, content: str) -> bytes:
    """Gera Excel com openpyxl."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Dados"

    # Tentar parsear content como JSON (array de objetos ou array de arrays)
    try:
        data = json.loads(content)
        if isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                # Array de objetos: headers das keys
                headers = list(data[0].keys())
                ws.append(headers)
                for row in data:
                    ws.append([str(row.get(h, "")) for h in headers])
            elif isinstance(data[0], list):
                # Array de arrays
                for row in data:
                    ws.append([str(cell) for cell in row])
            else:
                ws.append([title])
                ws.append([content])
        else:
            ws.append([title])
            ws.append([content])
    except (json.JSONDecodeError, TypeError):
        # Não é JSON, tratar como texto
        ws.append([title])
        for line in content.split("\n"):
            if line.strip():
                ws.append([line.strip()])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def generate_pptx(title: str, content: str) -> bytes:
    """Gera PPTX com python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Separar por marcador SLIDE:
    import re
    slides_content = re.split(r'SLIDE:', content, flags=re.IGNORECASE)
    slides_content = [s.strip() for s in slides_content if s.strip()]

    if not slides_content:
        # Fallback: um slide com tudo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    else:
        for slide_text in slides_content:
            lines = slide_text.split("\n")
            slide_title = lines[0].replace("*", "").replace("#", "").strip()
            slide_body = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_body

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _inject_export_css(html: str, vw: int) -> str:
    """Injeta CSS override para export: remove body flex/padding/overflow do design_contract."""
    css = (
        '<style id="arcco-export-override">'
        'html,body{margin:0!important;padding:0!important;'
        'overflow:visible!important;display:block!important;'
        f'width:{vw}px!important;height:auto!important;'
        'min-height:auto!important;align-items:initial!important;'
        'justify-content:initial!important}'
        '</style>'
    )
    if re.search(r"</head>", html, re.IGNORECASE):
        return re.sub(r"</head>", css + "</head>", html, count=1, flags=re.IGNORECASE)
    if re.search(r"<body", html, re.IGNORECASE):
        return re.sub(r"(<body)", css + r"\1", html, count=1, flags=re.IGNORECASE)
    return css + html


def _inject_tailwind_if_needed(html_content: str) -> str:
    """Injeta Tailwind CDN se o HTML não tiver estilos próprios."""
    tailwind_cdn = '<script src="https://cdn.tailwindcss.com"></script>'
    if "tailwindcss" not in html_content and "cdn.tailwindcss" not in html_content:
        if "<head>" in html_content:
            return html_content.replace("<head>", f"<head>{tailwind_cdn}", 1)
        elif "<html" in html_content:
            return tailwind_cdn + html_content
        return tailwind_cdn + html_content
    return html_content


async def html_to_screenshot(
    html_content: str,
    img_format: str = "png",
    slide_index: Optional[int] = None,
    slide_indices: Optional[list] = None,
    resolution: Optional[str] = None,
    canvas_preset: Optional[str] = None,
    viewport_width: Optional[int] = None,
    viewport_height: Optional[int] = None,
) -> Union[bytes, tuple]:
    """
    Captura screenshot de um HTML via Playwright.
    - slide_index: None = todos os slides (ZIP), int = slide específico
    - resolution: viewport size ("hd-720", "hd-1080")
    Retorna bytes (imagem única) ou tuple (zip_bytes, mime, ext) para multi-slide.
    Se viewport_width/viewport_height forem fornecidos, têm prioridade sobre resolution/canvas_preset.
    """
    inject = _inject_tailwind_if_needed(html_content)
    fmt = img_format.lower() if img_format.lower() in ("png", "jpeg") else "png"
    if viewport_width and viewport_height:
        vw, vh = viewport_width, viewport_height
    elif canvas_preset and canvas_preset in CANVAS_PRESETS:
        spec = CANVAS_PRESETS[canvas_preset]
        vw, vh = spec["width"], spec["height"]
    else:
        vw, vh = _RES_VIEWPORTS.get(resolution or "hd-720", (1280, 720))

    # Export CSS override: remove body flex/padding/overflow do design_contract
    inject = _inject_export_css(inject, vw)

    def _sync() -> Union[bytes, tuple]:
        import zipfile
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page(viewport={"width": vw, "height": vh})
            try:
                page.set_content(inject, wait_until="networkidle", timeout=30_000)
                slide_count = page.evaluate("() => document.querySelectorAll('.slide, .slide-container').length")

                if slide_count and int(slide_count) > 1:
                    slide_els = page.query_selector_all('.slide, .slide-container')
                    total = int(slide_count)

                    # Determina quais índices capturar
                    if slide_indices is not None:
                        capture_indices = sorted(set(i for i in slide_indices if 0 <= i < total))
                    elif slide_index is not None:
                        capture_indices = [slide_index] if 0 <= slide_index < total else []
                    else:
                        capture_indices = list(range(total))

                    if len(capture_indices) == 1:
                        shot = _capture_slide(page, slide_els, capture_indices[0], fmt)
                        return shot if shot else page.screenshot(full_page=False, type=fmt)
                    else:
                        # Múltiplos slides → ZIP
                        screenshots: list[bytes] = []
                        for i in capture_indices:
                            shot = _capture_slide(page, slide_els, i, fmt)
                            if shot:
                                screenshots.append(shot)

                        zip_buf = io.BytesIO()
                        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                            for idx, img_bytes in enumerate(screenshots):
                                zf.writestr(f"slide_{idx + 1}.{fmt}", img_bytes)
                        return (zip_buf.getvalue(), "application/zip", "zip")
                else:
                    # Documento simples
                    return page.screenshot(full_page=False, type=fmt)
            finally:
                browser.close()

    return await asyncio.to_thread(_sync)


async def html_to_pptx(
    html_content: str,
    title: str = "Apresentação",
    slide_index: Optional[int] = None,
    slide_indices: Optional[list] = None,
    page_size: Optional[str] = None,
    canvas_preset: Optional[str] = None,
    viewport_width: Optional[int] = None,
    viewport_height: Optional[int] = None,
) -> bytes:
    """
    Converte HTML com slides (.slide) em PPTX via screenshots Playwright.
    Cada <section class="slide"> vira um slide independente.
    Se não houver .slide, usa um único slide com screenshot full.
    O viewport e dimensões do PPTX se adaptam ao page_size escolhido.
    Se viewport_width/viewport_height forem fornecidos, têm prioridade.
    """
    inject = _inject_tailwind_if_needed(html_content)

    # Prioridade: viewport explícito > canvas_preset > page_size
    if viewport_width and viewport_height:
        vw, vh = viewport_width, viewport_height
        emu_per_px = 9525
        pptx_w, pptx_h = vw * emu_per_px, vh * emu_per_px
    elif canvas_preset and canvas_preset in CANVAS_PRESETS:
        spec = CANVAS_PRESETS[canvas_preset]
        vw, vh = spec["width"], spec["height"]
        emu_per_px = 9525
        pptx_w, pptx_h = vw * emu_per_px, vh * emu_per_px
    else:
        vw, vh = _PAGE_VIEWPORTS.get(page_size or "widescreen", (1920, 1080))
        # Dimensões do PPTX em EMUs (English Metric Units) — 914400 EMUs = 1 inch
        _PPTX_DIMS = {
            "widescreen":       (12192000, 6858000),   # 13.33 x 7.5 in
            "a4-landscape":     (10691812, 7562088),    # 11.69 x 8.27 in
            "a4-portrait":      (7562088, 10691812),    # 8.27 x 11.69 in
            "letter-landscape": (10058400, 7772400),    # 11.0 x 8.5 in
            "letter-portrait":  (7772400, 10058400),    # 8.5 x 11.0 in
        }
        pptx_w, pptx_h = _PPTX_DIMS.get(page_size or "widescreen", (12192000, 6858000))

    # Export CSS override: remove body flex/padding/overflow do design_contract
    inject = _inject_export_css(inject, vw)

    def _sync() -> bytes:
        import io as _io
        from playwright.sync_api import sync_playwright
        from pptx import Presentation as PptxPresentation

        screenshots: list[bytes] = []

        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page(viewport={"width": vw, "height": vh})
            try:
                page.set_content(inject, wait_until="networkidle", timeout=30_000)
                slide_count = page.evaluate("() => document.querySelectorAll('.slide, .slide-container').length")

                if slide_count and slide_count > 0:
                    # Determina quais slides capturar
                    if slide_indices is not None:
                        indices = sorted(set(i for i in slide_indices if 0 <= i < int(slide_count)))
                        if not indices:
                            indices = [0]
                    elif slide_index is not None:
                        indices = [slide_index] if 0 <= slide_index < slide_count else [0]
                    else:
                        indices = list(range(int(slide_count)))

                    slide_els = page.query_selector_all('.slide, .slide-container')
                    for i in indices:
                        shot = _capture_slide(page, slide_els, i, "png")
                        if shot:
                            screenshots.append(shot)
                else:
                    screenshots.append(page.screenshot(full_page=False, type="png"))
            finally:
                browser.close()

        # Monta o PPTX com dimensões dinâmicas
        prs = PptxPresentation()
        prs.slide_width = pptx_w
        prs.slide_height = pptx_h
        blank_layout = prs.slide_layouts[6]

        for img_bytes in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                _io.BytesIO(img_bytes), 0, 0,
                prs.slide_width, prs.slide_height
            )

        buf = _io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    return await asyncio.to_thread(_sync)


def _text_to_html(title: str, content: str) -> str:
    """Converte texto/markdown simples em HTML bonito para exportação via Playwright."""
    import html as html_lib
    lines = content.split("\n")
    body_parts = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            body_parts.append("<br/>")
        elif stripped.startswith("### "):
            body_parts.append(f'<h3 class="text-lg font-semibold text-gray-700 mt-4 mb-1">{html_lib.escape(stripped[4:])}</h3>')
        elif stripped.startswith("## "):
            body_parts.append(f'<h2 class="text-xl font-bold text-gray-800 mt-6 mb-2">{html_lib.escape(stripped[3:])}</h2>')
        elif stripped.startswith("# "):
            body_parts.append(f'<h1 class="text-2xl font-bold text-gray-900 mt-6 mb-3">{html_lib.escape(stripped[2:])}</h1>')
        elif stripped.startswith("- ") or stripped.startswith("* "):
            body_parts.append(f'<li class="ml-5 list-disc text-gray-700">{html_lib.escape(stripped[2:])}</li>')
        else:
            body_parts.append(f'<p class="text-gray-700 leading-relaxed my-1">{html_lib.escape(stripped)}</p>')

    body_html = "\n".join(body_parts)
    safe_title = html_lib.escape(title)
    return f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
  <meta charset="UTF-8"/>
  <script src="https://cdn.tailwindcss.com"></script>
  <title>{safe_title}</title>
</head>
<body class="bg-white font-sans p-12 max-w-3xl mx-auto">
  <h1 class="text-3xl font-bold text-gray-900 border-b-2 border-indigo-600 pb-4 mb-8">{safe_title}</h1>
  <div class="prose text-base">
    {body_html}
  </div>
</body>
</html>"""


# Mapa de tipos para geração
FILE_GENERATORS = {
    "pdf": {
        "func": generate_pdf,
        "mime": "application/pdf",
        "ext": "pdf",
    },
    "docx": {
        "func": generate_docx,
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ext": "docx",
    },
    "excel": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "xlsx": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "pptx": {
        "func": generate_pptx,
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ext": "pptx",
    },
}


async def generate_file(file_type: str, title: str, content: str) -> Tuple[str, str]:
    """
    Gera arquivo e faz upload ao Supabase.

    Returns:
        (url_download, mensagem)
    """
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    file_type = file_type.lower()

    if file_type not in FILE_GENERATORS:
        raise ValueError(f"Tipo de arquivo inválido: {file_type}. Suportados: {list(FILE_GENERATORS.keys())}")

    gen = FILE_GENERATORS[file_type]

    # Gerar arquivo
    file_bytes = gen["func"](title, content)

    # Filename seguro
    import re
    safe_title = re.sub(r'[^a-z0-9]', '_', title.lower())[:50]
    filename = f"{safe_title}.{gen['ext']}"

    # Upload ao Supabase
    url = upload_to_supabase(
        bucket=config.supabase_storage_bucket,
        filename=filename,
        file_content=file_bytes,
        content_type=gen["mime"],
    )

    return url, f"{file_type.upper()} gerado com sucesso."

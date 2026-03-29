"""
Execução de ferramentas para os agentes especialistas.

Cada função de ferramenta está isolada aqui, desacoplada do endpoint HTTP.
"""

import asyncio
import io
import logging
import mimetypes
import os
import tempfile
import time
from pathlib import Path

import httpx

logger = logging.getLogger(__name__)

_E2B_ARTIFACT_ROOT = "/home/user"
_E2B_ARTIFACT_SCAN_DEPTH = 8
_E2B_ARTIFACT_MAX_BYTES = 20 * 1024 * 1024
_E2B_ARTIFACT_EXCLUDED_DIRS = {
    "__pycache__",
    ".cache",
    ".config",
    ".local",
    ".npm",
    ".venv",
    "node_modules",
}


async def execute_tool(func_name: str, func_args: dict, user_id: str | None = None) -> str:
    """Despachante principal: executa a ferramenta e retorna resultado como string."""
    if func_name == "web_search":
        return await _web_search(func_args.get("query", ""))

    elif func_name == "web_fetch":
        return await _web_fetch(func_args.get("url", ""))

    elif func_name == "generate_pdf":
        return await _generate_pdf(func_args)

    elif func_name == "generate_excel":
        return await _generate_excel(func_args)

    elif func_name == "execute_python":
        return await _execute_python(func_args.get("code", ""))

    elif func_name == "fetch_file_content":
        return await _fetch_file_content(func_args.get("url", ""))

    elif func_name == "read_session_file":
        return await _read_session_file(
            func_args.get("session_id", ""),
            func_args.get("file_name", ""),
            func_args.get("query"),
        )

    elif func_name == "modify_excel":
        return await _modify_excel(func_args)

    elif func_name == "modify_pptx":
        return await _modify_pptx(func_args)

    elif func_name == "modify_pdf":
        return await _modify_pdf(func_args)

    elif func_name == "ask_browser":
        return await _ask_browser(func_args)

    elif func_name == "generate_pdf_template":
        return await _generate_pdf_template(func_args)

    # ── Arcco Computer tools ──
    elif func_name == "list_computer_files":
        return await _list_computer_files(func_args, user_id)

    elif func_name == "read_computer_file":
        return await _read_computer_file(func_args, user_id)

    elif func_name == "manage_computer_file":
        return await _manage_computer_file(func_args, user_id)

    # ── Spy Pages (SimilarWeb via Apify) ──
    elif func_name == "analyze_web_pages":
        return await _analyze_web_pages(func_args)

    return f"Ferramenta desconhecida: {func_name}"


# ── Implementações ─────────────────────────────────────────────────────────────

async def _web_search(query: str) -> str:
    from backend.services.search_service import search_web_formatted
    return await search_web_formatted(query)


async def _web_fetch(url: str) -> str:
    try:
        from bs4 import BeautifulSoup
        async with httpx.AsyncClient(timeout=20.0) as client:
            response = await client.get(
                url, headers={"User-Agent": "ArccoAgent/2.0"}, follow_redirects=True
            )
            html = response.text

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form", "svg", "noscript"]):
            tag.decompose()

        text = soup.get_text(separator=" ", strip=True)
        title = soup.title.string if soup.title else url

        if len(text) > 20000:
            text = text[:20000] + "... [Truncado]"

        return f"**Conteúdo de {url}**\n**Título:** {title}\n\n{text}"
    except Exception as e:
        return f"Erro ao ler URL ({url}): {e}"


async def _ask_browser(args: dict) -> str:
    """
    Navega remotamente via Browserbase + Playwright (CDP).
    Delega toda a lógica para browser_service.execute_browserbase_task.
    """
    from backend.services.browser_service import execute_browserbase_task

    url = args.get("url", "")
    if not url:
        logger.error("[BROWSER] URL não fornecida para o Browser Agent. Args: %s", args)
        return "Erro: URL não fornecida para o Browser Agent."

    try:
        return await execute_browserbase_task(
            url=url,
            actions=args.get("actions", []),
            wait_for=int(args.get("wait_for") or 0),
            mobile=bool(args.get("mobile", False)),
            include_tags=args.get("include_tags"),
            exclude_tags=args.get("exclude_tags"),
        )
    except Exception as exc:
        logger.exception("[BROWSER] Falha ao executar Browser Agent em %s", url)
        return f"Erro ao executar Browser Agent: {exc}"


async def _generate_pdf(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    html_content = (args.get("html_content") or "").strip()

    if html_content:
        # Modo Playwright: HTML+Tailwind → PDF de alta qualidade
        from backend.services.file_service import generate_pdf_playwright
        logger.info("[PDF] Modo Playwright (HTML+Tailwind)")
        pdf_bytes = await generate_pdf_playwright(html_content)
    else:
        # Modo texto: reportlab (fallback)
        from backend.services.file_service import generate_pdf
        logger.info("[PDF] Modo reportlab (texto)")
        title = args.get("title", "documento")
        content = args.get("content", "")
        pdf_bytes = await asyncio.to_thread(generate_pdf, title, content)

    filename = args.get("filename", f"doc-{int(time.time())}")
    if not filename.endswith(".pdf"):
        filename += ".pdf"

    url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        pdf_bytes,
        "application/pdf",
    )
    return (
        f"PDF gerado com sucesso. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF]({url})"
    )


async def _generate_pdf_template(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase
    from backend.services.file_service import generate_pdf_from_template

    config = get_config()
    template_name = args.get("template_name", "relatorio")
    data = args.get("data", {})

    try:
        logger.info(f"[PDF] Gerando via template Jinja2: '{template_name}'")
        pdf_bytes = await generate_pdf_from_template(template_name, data)
    except Exception as e:
        return f"Erro ao gerar PDF com template '{template_name}': {e}"

    filename = args.get("filename", f"{template_name}-{int(time.time())}")
    if not filename.endswith(".pdf"):
        filename += ".pdf"

    url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        pdf_bytes,
        "application/pdf",
    )
    return (
        f"PDF gerado via template '{template_name}'. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF]({url})"
    )


async def _generate_excel(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase
    from openpyxl import Workbook

    config = get_config()

    def sync_excel():
        headers = [str(h) for h in args.get("headers", [])]
        rows = [[str(c) for c in row] for row in args.get("rows", [])]
        title = args.get("title", "Planilha")[:31]

        wb = Workbook()
        ws = wb.active
        ws.title = title
        ws.append(headers)
        for row in rows:
            ws.append(row)

        buffer = io.BytesIO()
        wb.save(buffer)
        file_bytes = buffer.getvalue()

        filename = args.get("filename", f"planilha-{int(time.time())}")
        if not filename.endswith(".xlsx"):
            filename += ".xlsx"

        return upload_to_supabase(
            config.supabase_storage_bucket,
            filename,
            file_bytes,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    url = await asyncio.to_thread(sync_excel)
    return (
        f"Planilha Excel gerada. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Planilha]({url})"
    )


async def _execute_python(code: str) -> str:
    try:
        from e2b_code_interpreter import Sandbox
    except ImportError:
        logger.exception("[E2B] Pacote e2b_code_interpreter não está instalado.")
        return (
            "Erro: pacote `e2b-code-interpreter` não está instalado. "
            "Execute: pip install e2b-code-interpreter"
        )

    from backend.core.config import get_config

    config = get_config()

    # Sempre busca a chave E2B fresca do Supabase (ignora o singleton cacheado no startup).
    # Isso garante que trocar a chave no Supabase Admin tem efeito imediato, sem restart.
    def _fetch_e2b_key_from_supabase() -> str | None:
        try:
            from backend.core.supabase_client import get_supabase_client
            db = get_supabase_client()
            for provider in ("e2b", "e2b_api_key"):
                rows = db.query("ApiKeys", "api_key", {"provider": provider})
                if rows and rows[0].get("api_key"):
                    return rows[0]["api_key"]
        except Exception as exc:
            logger.warning("[E2B] Falha ao buscar chave do Supabase: %s", exc)
        return None

    e2b_api_key = await asyncio.to_thread(_fetch_e2b_key_from_supabase)
    if not e2b_api_key:
        # Fallback: env var ou config cacheado no startup
        e2b_api_key = os.getenv("E2B_API_KEY") or config.e2b_api_key

    if not e2b_api_key:
        logger.error("[E2B] Variável E2B_API_KEY ausente; sandbox indisponível.")
        return (
            "Erro: A variável E2B_API_KEY não está configurada. "
            "Adicione provider='e2b' ou provider='e2b_api_key' na tabela ApiKeys do Supabase, ou defina E2B_API_KEY no ambiente."
        )

    logger.info("[E2B] Usando chave: %s...", e2b_api_key[:15])

    if not code or not code.strip():
        logger.error("[E2B] Tentativa de executar código vazio.")
        return "Erro: nenhum código Python foi fornecido para execução."

    try:
        from backend.core.supabase_client import upload_to_supabase

        def snapshot_artifacts(sandbox) -> dict[str, object]:
            entries = sandbox.files.list(_E2B_ARTIFACT_ROOT, depth=_E2B_ARTIFACT_SCAN_DEPTH)
            snapshot: dict[str, object] = {}
            for entry in entries:
                entry_type = getattr(getattr(entry, "type", None), "value", getattr(entry, "type", None))
                if entry_type != "file":
                    continue

                entry_path = str(getattr(entry, "path", "") or "")
                if not entry_path.startswith(f"{_E2B_ARTIFACT_ROOT}/"):
                    continue

                rel_path = entry_path.removeprefix(f"{_E2B_ARTIFACT_ROOT}/")
                parts = [part for part in Path(rel_path).parts if part not in ("", ".")]
                if not parts:
                    continue
                if any(part.startswith(".") for part in parts):
                    continue
                if any(part in _E2B_ARTIFACT_EXCLUDED_DIRS for part in parts[:-1]):
                    continue

                snapshot[entry_path] = entry
            return snapshot

        def upload_generated_artifacts(sandbox, before_snapshot: dict[str, object]) -> list[tuple[str, str]]:
            after_snapshot = snapshot_artifacts(sandbox)
            uploaded: list[tuple[str, str]] = []

            for path, entry in sorted(after_snapshot.items()):
                previous = before_snapshot.get(path)
                if previous is not None:
                    same_size = getattr(previous, "size", None) == getattr(entry, "size", None)
                    same_mtime = getattr(previous, "modified_time", None) == getattr(entry, "modified_time", None)
                    if same_size and same_mtime:
                        continue

                size = int(getattr(entry, "size", 0) or 0)
                if size <= 0 or size > _E2B_ARTIFACT_MAX_BYTES:
                    logger.info("[E2B] Ignorando artefato %s (size=%s).", path, size)
                    continue

                rel_path = Path(path).relative_to(_E2B_ARTIFACT_ROOT)
                safe_rel = str(rel_path).replace(os.sep, "_")
                filename = rel_path.name
                content_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"

                try:
                    file_bytes = bytes(sandbox.files.read(path, format="bytes"))
                    storage_name = f"python-artifacts/{safe_rel}"
                    public_url = upload_to_supabase(
                        config.supabase_storage_bucket,
                        storage_name,
                        file_bytes,
                        content_type,
                    )
                    uploaded.append((filename, public_url))
                    logger.info("[E2B] Artefato enviado ao Supabase: %s -> %s", path, public_url)
                except Exception as exc:
                    logger.exception("[E2B] Falha ao enviar artefato gerado '%s': %s", path, exc)

            return uploaded

        def _run_code_and_collect(sandbox, current_code: str, before_snapshot) -> tuple[str, str | None]:
            """Executa codigo no sandbox e retorna (resultado, erro_ou_none)."""
            execution_result = sandbox.run_code(current_code)

            parts = []
            error_msg = None

            # stdout de print() — vai para logs.stdout na v2
            if execution_result.logs and execution_result.logs.stdout:
                parts.append("".join(execution_result.logs.stdout))

            # Resultado de expressão (ex: 2+2 sem print) — vai para .text
            if execution_result.text:
                parts.append(execution_result.text)

            # Erros
            if execution_result.error:
                error_msg = f"{execution_result.error.name}: {execution_result.error.value}"
                parts.append(f"ERRO: {error_msg}")

            # stderr — pode indicar erro mesmo sem execution_result.error
            if execution_result.logs and execution_result.logs.stderr:
                stderr_text = "".join(execution_result.logs.stderr)
                parts.append("STDERR: " + stderr_text)
                if not error_msg and ("Error" in stderr_text or "Traceback" in stderr_text):
                    error_msg = stderr_text[-500:]

            # Resultados extras (mídia, display)
            if execution_result.results:
                for media in execution_result.results:
                    if hasattr(media, 'text') and media.text:
                        parts.append(f"RESULT: {media.text}")

            uploaded_artifacts = upload_generated_artifacts(sandbox, before_snapshot)
            if uploaded_artifacts:
                parts.append("Arquivos gerados:")
                parts.extend(f"[{filename}]({url})" for filename, url in uploaded_artifacts)

            result_str = "\n".join(parts).strip()
            return result_str, error_msg

        def sync_e2b_exec():
            logger.info("[E2B] Iniciando execução de código Python (%s chars).", len(code))
            sandbox = Sandbox(api_key=e2b_api_key)
            try:
                before_snapshot = snapshot_artifacts(sandbox)
                result_str, error_msg = _run_code_and_collect(sandbox, code, before_snapshot)

                if not error_msg:
                    return result_str if result_str else "(Código executado com sucesso sem output. Use print() para ver resultados.)"

                # ── Auto-healing: marca codigo e erro para correcao ──
                return f"__AUTOFIX__|{error_msg}|__CODE__|{code}|__RESULT__|{result_str}"
            finally:
                sandbox.kill()

        first_result = await asyncio.to_thread(sync_e2b_exec)

        # Se nao precisa de auto-fix, retorna direto
        if not first_result.startswith("__AUTOFIX__"):
            logger.info("[E2B] Execução concluída com %s chars de saída.", len(first_result))
            return first_result

        # ── Loop de Auto-Healing (max 2 retries) ──────────────────────────────
        _MAX_AUTOFIX = 2
        current_code = code
        last_error = first_result.split("|__CODE__|")[0].replace("__AUTOFIX__|", "")
        last_result = first_result.split("|__RESULT__|")[-1] if "|__RESULT__|" in first_result else ""

        for attempt in range(1, _MAX_AUTOFIX + 1):
            logger.info("[E2B-AUTOFIX] Tentativa %d/%d de auto-correção...", attempt, _MAX_AUTOFIX)

            # Pede ao LLM para corrigir
            try:
                from backend.core.llm import call_openrouter
                fix_response = await call_openrouter(
                    messages=[{
                        "role": "user",
                        "content": (
                            "Voce e um expert em Python. O seguinte codigo falhou com este erro:\n\n"
                            f"CODIGO:\n```python\n{current_code}\n```\n\n"
                            f"ERRO:\n{last_error}\n\n"
                            "Reescreva o codigo completo para corrigir o problema. "
                            "Retorne APENAS o codigo Python puro, sem formatacao markdown, "
                            "sem ```python, sem explicacao."
                        ),
                    }],
                    model="openai/gpt-4o-mini",
                    max_tokens=4000,
                    temperature=0.1,
                )
                fixed_code = fix_response["choices"][0]["message"]["content"].strip()
                # Remove blocos markdown se o LLM ignorar a instrucao
                if fixed_code.startswith("```"):
                    import re
                    fixed_code = re.sub(r"^```(?:python)?\s*\n?", "", fixed_code)
                    fixed_code = re.sub(r"\n?```\s*$", "", fixed_code)
                    fixed_code = fixed_code.strip()

                if not fixed_code or len(fixed_code) < 10:
                    logger.warning("[E2B-AUTOFIX] LLM retornou código vazio/curto, abortando.")
                    break

                current_code = fixed_code
            except Exception as fix_exc:
                logger.error("[E2B-AUTOFIX] Falha ao chamar LLM para correção: %s", fix_exc)
                break

            # Re-executa o codigo corrigido
            def sync_retry(retry_code=current_code):
                sandbox = Sandbox(api_key=e2b_api_key)
                try:
                    before = snapshot_artifacts(sandbox)
                    return _run_code_and_collect(sandbox, retry_code, before)
                finally:
                    sandbox.kill()

            try:
                result_str, error_msg = await asyncio.to_thread(sync_retry)
            except Exception as retry_exc:
                logger.error("[E2B-AUTOFIX] Falha na re-execução: %s", retry_exc)
                last_error = str(retry_exc)
                continue

            if not error_msg:
                logger.info("[E2B-AUTOFIX] Código corrigido com sucesso na tentativa %d!", attempt)
                return result_str if result_str else "(Código executado com sucesso sem output. Use print() para ver resultados.)"

            last_error = error_msg
            last_result = result_str
            logger.warning("[E2B-AUTOFIX] Tentativa %d ainda com erro: %s", attempt, error_msg[:100])

        # Todas as tentativas falharam
        logger.error("[E2B-AUTOFIX] Falha após %d tentativas de auto-correção.", _MAX_AUTOFIX)
        return last_result if last_result else f"Falha após {_MAX_AUTOFIX + 1} tentativas de execução. Último erro: {last_error}"
    except Exception as e:
        logger.exception("[E2B] Erro durante a execução do sandbox Python.")
        return f"Erro na execução (E2B Sandbox): {e}"


# ── Modificador de Arquivos ────────────────────────────────────────────────────

async def _fetch_file_content(url: str) -> str:
    """Baixa um arquivo e retorna sua estrutura como texto legível."""
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar arquivo: HTTP {response.status_code}"

        content_type = response.headers.get("content-type", "").lower()
        file_bytes = response.content
        url_lower = url.lower().split("?")[0]  # ignora query params

        if "spreadsheet" in content_type or url_lower.endswith(".xlsx"):
            return await asyncio.to_thread(_read_excel_structure, file_bytes)
        elif "presentation" in content_type or url_lower.endswith(".pptx"):
            return await asyncio.to_thread(_read_pptx_structure, file_bytes)
        elif "pdf" in content_type or url_lower.endswith(".pdf"):
            return await asyncio.to_thread(_read_pdf_text, file_bytes)
        else:
            return f"Tipo de arquivo não identificado (content-type: {content_type}). URL: {url}"
    except Exception as e:
        return f"Erro ao ler arquivo: {e}"


async def _read_session_file(
    session_id: str,
    file_name: str,
    query: str | None = None,
) -> str:
    from backend.services.ephemeral_rag_service import (
        format_chunk_results,
        search_relevant_chunks,
    )
    from backend.services.session_file_service import (
        SessionFileError,
        SessionFileNotFoundError,
        get_session_file_by_name,
        touch_session,
    )

    try:
        if not session_id:
            return "Erro: session_id não informado para leitura do arquivo da sessão."
        if not file_name:
            return "Erro: file_name não informado para leitura do arquivo da sessão."

        touch_session(session_id)
        entry = get_session_file_by_name(session_id, file_name)
        status = entry.get("status", "uploaded")
        extracted_path = entry.get("extracted_text_path", "")

        if status in {"uploaded", "processing"}:
            return (
                f"O arquivo '{entry.get('original_name', file_name)}' ainda está em processamento. "
                "Avise o usuário que o OCR/leitura ainda está rodando e peça para tentar novamente em instantes."
            )

        if status == "failed":
            error_message = entry.get("error") or "Falha desconhecida durante a extração."
            return (
                f"Não foi possível ler o arquivo '{entry.get('original_name', file_name)}'. "
                f"Erro registrado: {error_message}"
            )

        if not extracted_path:
            return f"Erro: o arquivo '{entry.get('original_name', file_name)}' não possui texto extraído registrado."

        text = await asyncio.to_thread(
            lambda: open(extracted_path, "r", encoding="utf-8").read()
        )
        if not text.strip():
            return f"O arquivo '{entry.get('original_name', file_name)}' foi processado, mas não contém texto legível."

        if query:
            chunks = search_relevant_chunks(text, query)
            return format_chunk_results(entry.get("original_name", file_name), chunks)

        preview = text[:3000].strip()
        if len(text) > 3000:
            preview += "\n\n... [conteúdo truncado]"
        return f"Prévia do arquivo '{entry.get('original_name', file_name)}':\n\n{preview}"
    except SessionFileNotFoundError as exc:
        return str(exc)
    except SessionFileError as exc:
        return f"Erro ao consultar arquivos da sessão: {exc}"
    except Exception as exc:
        logger.error(
            "Erro inesperado ao ler arquivo de sessão '%s' (%s): %s",
            file_name,
            session_id,
            exc,
        )
        return f"Erro ao ler arquivo da sessão: {exc}"


def _read_excel_structure(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes))
    lines = [f"Planilha Excel — {len(wb.sheetnames)} aba(s): {', '.join(wb.sheetnames)}"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append(f"\nAba '{sheet_name}': vazia")
            continue
        headers = [str(c) if c is not None else "" for c in rows[0]]
        lines.append(f"\nAba '{sheet_name}' — {ws.max_row} linha(s), {ws.max_column} coluna(s)")
        lines.append(f"Cabeçalhos (linha 1): {headers}")
        for i, row in enumerate(rows[1:6], start=2):
            lines.append(f"  Linha {i}: {[str(c) if c is not None else '' for c in row]}")
        if ws.max_row > 6:
            lines.append(f"  ... ({ws.max_row - 6} linha(s) restante(s) não exibidas)")
    return "\n".join(lines)


def _read_pptx_structure(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = [f"Apresentação PPTX — {len(prs.slides)} slide(s)"]
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        preview = " | ".join(texts[:4]) if texts else "(sem texto)"
        lines.append(f"  Slide {i + 1}: {preview}")
    return "\n".join(lines)


def _read_pdf_text(file_bytes: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    lines = [f"PDF — {len(reader.pages)} página(s)"]
    total_chars = 0
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        snippet = text[:800].strip()
        lines.append(f"\n--- Página {i + 1} ---\n{snippet}")
        total_chars += len(text)
        if total_chars > 4000:
            lines.append("... [restante truncado]")
            break
    return "\n".join(lines)


async def _modify_excel(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    cell_updates = args.get("cell_updates", [])
    append_rows = args.get("append_rows", [])
    output_filename = args.get("output_filename", f"planilha-modificada")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar planilha: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar planilha: {e}"

    def sync_modify():
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes))

        for update in cell_updates:
            sheet_name = update.get("sheet", "")
            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws[update["cell"]] = update["value"]

        for row_def in append_rows:
            sheet_name = row_def.get("sheet", "")
            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws.append(row_def.get("values", []))

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".xlsx") else output_filename + ".xlsx"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    return (
        f"Planilha modificada com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Planilha Modificada]({upload_url})"
    )


async def _modify_pptx(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    text_replacements = args.get("text_replacements", [])
    output_filename = args.get("output_filename", f"apresentacao-modificada")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar apresentação: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar apresentação: {e}"

    def sync_modify():
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for rep in text_replacements:
                            if rep["find"] in run.text:
                                run.text = run.text.replace(rep["find"], rep["replace"])

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".pptx") else output_filename + ".pptx"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return (
        f"Apresentação modificada com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Apresentação Modificada]({upload_url})"
    )


async def _modify_pdf(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    text_replacements = args.get("text_replacements", [])
    append_content = args.get("append_content", "")
    output_filename = args.get("output_filename", f"documento-modificado")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar PDF: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar PDF: {e}"

    def sync_modify():
        import PyPDF2
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        # Extrai texto de todas as páginas
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        full_text = "\n\n".join(page.extract_text() or "" for page in reader.pages)

        # Aplica substituições
        for rep in text_replacements:
            full_text = full_text.replace(rep["find"], rep["replace"])

        # Adiciona conteúdo extra
        if append_content:
            full_text += f"\n\n{append_content}"

        # Regera o PDF com reportlab
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        for line in full_text.split("\n"):
            if line.strip():
                safe_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe_line, styles["Normal"]))
                story.append(Spacer(1, 4))
        doc.build(story)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".pdf") else output_filename + ".pdf"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/pdf",
    )
    return (
        f"PDF modificado com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF Modificado]({upload_url})"
    )


# ── Arcco Computer Tools ─────────────────────────────────────────────────────


async def _list_computer_files(args: dict, user_id: str | None) -> str:
    """Lista arquivos do usuário no Arcco Computer."""
    if not user_id:
        return "Erro: user_id não disponível. O usuário precisa estar logado."

    from backend.core.supabase_client import get_supabase_client

    folder_path = args.get("folder_path", "/")
    client = get_supabase_client()

    try:
        files = client.query(
            "user_files",
            select="id,file_name,file_type,size_bytes,folder_path,created_at",
            filters={"user_id": user_id, "folder_path": folder_path},
            order="created_at.desc",
        )

        if not files:
            return f"Nenhum arquivo encontrado na pasta '{folder_path}'."

        # Filtra markers de pasta
        real_files = [f for f in files if f.get("file_type") != "folder"]

        # Busca subpastas
        all_paths = client.query(
            "user_files",
            select="folder_path",
            filters={"user_id": user_id},
        )
        prefix = folder_path if folder_path == "/" else folder_path + "/"
        subfolders = set()
        for row in all_paths:
            fp = row.get("folder_path", "")
            if fp and fp != folder_path and fp.startswith(prefix):
                rest = fp[len(prefix):]
                next_level = rest.split("/")[0]
                if next_level:
                    subfolders.add(next_level)

        lines = [f"Pasta: {folder_path}"]
        if subfolders:
            lines.append(f"Subpastas: {', '.join(sorted(subfolders))}")
        lines.append(f"Arquivos ({len(real_files)}):")
        for f in real_files:
            size = f.get("size_bytes", 0)
            size_str = f"{size} B" if size < 1024 else f"{size / 1024:.1f} KB" if size < 1048576 else f"{size / 1048576:.1f} MB"
            lines.append(f"  - id={f['id']} | {f['file_name']} | {f.get('file_type', '?')} | {size_str}")

        return "\n".join(lines)
    except Exception as exc:
        logger.error("Erro ao listar arquivos do Computer: %s", exc)
        return f"Erro ao listar arquivos: {exc}"


async def _read_computer_file(args: dict, user_id: str | None) -> str:
    """Lê conteúdo de um arquivo do Arcco Computer."""
    if not user_id:
        return "Erro: user_id não disponível."

    from backend.core.supabase_client import get_supabase_client

    file_id = args.get("file_id", "")
    query = args.get("query")

    if not file_id:
        return "Erro: file_id é obrigatório."

    client = get_supabase_client()

    try:
        rows = client.query(
            "user_files",
            filters={"id": file_id, "user_id": user_id},
        )
        if not rows:
            return f"Arquivo não encontrado (id={file_id})."

        file_meta = rows[0]
        file_url = file_meta.get("file_url", "")
        file_name = file_meta.get("file_name", "arquivo")
        file_type = file_meta.get("file_type", "")

        if not file_url:
            return f"Arquivo '{file_name}' não possui URL."

        # Baixa o conteúdo
        async with httpx.AsyncClient(timeout=30.0) as http:
            resp = await http.get(file_url, follow_redirects=True)
        if resp.status_code != 200:
            return f"Erro ao baixar '{file_name}': HTTP {resp.status_code}"

        content_bytes = resp.content

        # Extrai texto baseado no tipo
        text = await asyncio.to_thread(_extract_text_from_bytes, content_bytes, file_name, file_type)

        if not text or not text.strip():
            return f"O arquivo '{file_name}' não contém texto legível."

        # RAG se query fornecida
        if query:
            try:
                from backend.services.ephemeral_rag_service import search_relevant_chunks, format_chunk_results
                chunks = search_relevant_chunks(text, query)
                return format_chunk_results(file_name, chunks)
            except Exception:
                pass  # fallback para preview simples

        preview = text[:4000].strip()
        if len(text) > 4000:
            preview += "\n\n... [conteúdo truncado — use query para buscar trechos específicos]"
        return f"Conteúdo de '{file_name}':\n\n{preview}"

    except Exception as exc:
        logger.error("Erro ao ler arquivo do Computer (id=%s): %s", file_id, exc)
        return f"Erro ao ler arquivo: {exc}"


def _extract_text_from_bytes(content: bytes, filename: str, mime_type: str) -> str:
    """Extrai texto de bytes baseado no tipo do arquivo."""
    name_lower = filename.lower()

    # PDF
    if "pdf" in mime_type or name_lower.endswith(".pdf"):
        return _read_pdf_text(content)

    # Excel
    if "spreadsheet" in mime_type or "excel" in mime_type or name_lower.endswith((".xlsx", ".xls")):
        return _read_excel_structure(content)

    # PPTX
    if "presentation" in mime_type or name_lower.endswith((".pptx", ".ppt")):
        return _read_pptx_structure(content)

    # DOCX
    if "document" in mime_type or name_lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return content.decode("utf-8", errors="replace")

    # Texto puro / CSV / JSON / MD
    if any(t in mime_type for t in ("text", "json", "csv", "markdown")) or name_lower.endswith(
        (".txt", ".csv", ".md", ".json", ".html", ".css", ".js")
    ):
        return content.decode("utf-8", errors="replace")

    # Fallback: tenta como texto
    try:
        return content.decode("utf-8", errors="replace")
    except Exception:
        return ""


async def _manage_computer_file(args: dict, user_id: str | None) -> str:
    """Gerencia arquivos no Arcco Computer: mover, renomear, criar pasta, salvar novo."""
    if not user_id:
        return "Erro: user_id não disponível."

    from backend.core.supabase_client import get_supabase_client, upload_to_supabase

    action = args.get("action", "")
    client = get_supabase_client()

    try:
        if action == "move":
            file_id = args.get("file_id", "")
            target = args.get("target_folder", "/")
            if not file_id:
                return "Erro: file_id é obrigatório para mover."
            client.update("user_files", {"folder_path": target}, {"id": file_id, "user_id": user_id})
            return f"Arquivo movido para '{target}'."

        elif action == "rename":
            file_id = args.get("file_id", "")
            new_name = args.get("new_name", "")
            if not file_id or not new_name:
                return "Erro: file_id e new_name são obrigatórios para renomear."
            client.update("user_files", {"file_name": new_name}, {"id": file_id, "user_id": user_id})
            return f"Arquivo renomeado para '{new_name}'."

        elif action == "create_folder":
            target = args.get("target_folder", "")
            if not target or target == "/":
                return "Erro: informe o caminho da pasta (ex: /Marketing)."
            client.insert("user_files", {
                "user_id": user_id,
                "file_name": ".folder",
                "file_url": "",
                "file_type": "folder",
                "folder_path": target,
                "size_bytes": 0,
            })
            return f"Pasta '{target}' criada."

        elif action == "save_new":
            file_name = args.get("file_name", f"arquivo_{int(time.time())}.txt")
            content = args.get("content", "")
            file_type = args.get("file_type", "text/plain")
            target_folder = args.get("target_folder", "/")

            if not content:
                return "Erro: content é obrigatório para salvar novo arquivo."

            content_bytes = content.encode("utf-8")
            storage_path = f"{user_id}/{int(time.time())}_{file_name}"

            sb_client = get_supabase_client()
            public_url = sb_client.storage_upload("user_artifacts", storage_path, content_bytes, file_type)

            client.insert("user_files", {
                "user_id": user_id,
                "file_name": file_name,
                "file_url": public_url,
                "file_type": file_type,
                "size_bytes": len(content_bytes),
                "folder_path": target_folder,
            })
            return f"Arquivo '{file_name}' salvo em '{target_folder}'."

        else:
            return f"Ação desconhecida: '{action}'. Use: move, rename, create_folder ou save_new."

    except Exception as exc:
        logger.error(f"[COMPUTER] Erro em manage_computer_file action='{action}': {exc}")
        return f"Erro ao executar ação '{action}': {exc}"


# ── Spy Pages (SimilarWeb via Apify) ─────────────────────────────────────────

async def _analyze_web_pages(args: dict) -> str:
    import json as _json
    from backend.services.apify_service import analyze_pages

    urls = args.get("urls", [])[:4]
    if not urls:
        return "Erro: nenhuma URL fornecida para análise."

    try:
        results = await analyze_pages(urls)
        return _json.dumps(results, ensure_ascii=False)
    except Exception as exc:
        logger.error(f"[EXECUTOR] Erro em analyze_web_pages: {exc}")
        return f"Erro ao analisar sites: {exc}"

    except Exception as exc:
        logger.error("Erro ao gerenciar arquivo do Computer (action=%s): %s", action, exc)
        return f"Erro ao executar '{action}': {exc}"
